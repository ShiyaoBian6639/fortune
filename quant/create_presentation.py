"""
Create PowerPoint presentation for Trump China Visit Stock Market Analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path

# Alias for convenience
RgbColor = RGBColor

PLOTS_DIR = Path(__file__).parent.parent / 'plots'
OUTPUT_DIR = Path(__file__).parent.parent


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = subtitle
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RgbColor(100, 100, 100)
    subtitle_para.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullet_points=None, image_path=None, image_width=None):
    """Add a content slide with optional image."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(0, 51, 102)

    # Add underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.9), Inches(9.4), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RgbColor(0, 102, 204)
    line.line.fill.background()

    content_top = 1.1

    if bullet_points:
        text_box = slide.shapes.add_textbox(Inches(0.3), Inches(content_top), Inches(4.5), Inches(5.5))
        tf = text_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = f"• {point}"
            p.font.size = Pt(14)
            p.space_after = Pt(8)
            p.font.color.rgb = RgbColor(50, 50, 50)

        if image_path:
            img_width = image_width or Inches(5)
            slide.shapes.add_picture(str(image_path), Inches(4.8), Inches(content_top), width=img_width)
    elif image_path:
        img_width = image_width or Inches(9)
        slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(content_top), width=img_width)

    return slide


def add_two_image_slide(prs, title, image1_path, image2_path, caption1="", caption2=""):
    """Add a slide with two images side by side."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(0, 51, 102)

    # Add images
    slide.shapes.add_picture(str(image1_path), Inches(0.2), Inches(1.2), width=Inches(4.8))
    slide.shapes.add_picture(str(image2_path), Inches(5), Inches(1.2), width=Inches(4.8))

    # Add captions
    if caption1:
        cap1_box = slide.shapes.add_textbox(Inches(0.2), Inches(6.8), Inches(4.8), Inches(0.5))
        cap1_frame = cap1_box.text_frame
        cap1_para = cap1_frame.paragraphs[0]
        cap1_para.text = caption1
        cap1_para.font.size = Pt(12)
        cap1_para.alignment = PP_ALIGN.CENTER

    if caption2:
        cap2_box = slide.shapes.add_textbox(Inches(5), Inches(6.8), Inches(4.8), Inches(0.5))
        cap2_frame = cap2_box.text_frame
        cap2_para = cap2_frame.paragraphs[0]
        cap2_para.text = caption2
        cap2_para.font.size = Pt(12)
        cap2_para.alignment = PP_ALIGN.CENTER

    return slide


def add_key_findings_slide(prs):
    """Add key findings slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Key Findings"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(0, 51, 102)

    # Findings boxes
    findings = [
        ("61.7%", "of stocks gained during the visit", "green"),
        ("+1.17%", "average market gain", "blue"),
        ("+33%", "top performer gain (雅克科技)", "orange"),
        ("$250B", "total trade deals signed", "purple"),
    ]

    for i, (number, desc, color) in enumerate(findings):
        left = 0.3 + (i * 2.4)

        # Number box
        colors = {
            "green": RgbColor(34, 139, 34),
            "blue": RgbColor(30, 144, 255),
            "orange": RgbColor(255, 140, 0),
            "purple": RgbColor(138, 43, 226),
        }

        num_box = slide.shapes.add_textbox(Inches(left), Inches(1.5), Inches(2.2), Inches(1))
        num_frame = num_box.text_frame
        num_para = num_frame.paragraphs[0]
        num_para.text = number
        num_para.font.size = Pt(40)
        num_para.font.bold = True
        num_para.font.color.rgb = colors[color]
        num_para.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(left), Inches(2.5), Inches(2.2), Inches(1))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_para = desc_frame.paragraphs[0]
        desc_para.text = desc
        desc_para.font.size = Pt(14)
        desc_para.alignment = PP_ALIGN.CENTER

    return slide


def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Trump's China Visit:\nStock Market Impact Analysis",
        "November 8-10, 2017 | $250 Billion Trade Deals"
    )

    # Slide 2: Executive Summary
    add_content_slide(
        prs,
        "Executive Summary",
        bullet_points=[
            "President Trump visited China on November 8-10, 2017",
            "First state visit by a U.S. president since the founding of PRC",
            "Over $250 billion in trade deals announced",
            "Major sectors impacted: Energy, Aviation, Semiconductors, Manufacturing",
            "Chinese stock market showed positive reaction with 61.7% of stocks gaining",
            "Average market gain of +1.17% during the 3-day visit period"
        ]
    )

    # Slide 3: The Visit and Trade Deals
    add_content_slide(
        prs,
        "The Historic Visit & Trade Deals",
        bullet_points=[
            "ENERGY: $84B West Virginia shale gas investment",
            "         $43B Alaska LNG deal with Sinopec",
            "",
            "AVIATION: $37B Boeing deal for 300 aircraft",
            "",
            "TECHNOLOGY: $12B Qualcomm deals with Chinese smartphone makers",
            "",
            "AUTOMOTIVE: $10B Ford vehicle exports",
            "                   $756M Ford-Zotye EV joint venture",
            "",
            "AGRICULTURE: $5B soybean purchase agreements",
            "",
            "FINANCE: $5B Goldman Sachs-CIC investment fund"
        ]
    )

    # Slide 4: Market Overview
    add_content_slide(
        prs,
        "Overall Market Reaction",
        image_path=PLOTS_DIR / 'market_overview.png',
        image_width=Inches(9)
    )

    # Slide 5: Distribution
    add_content_slide(
        prs,
        "Price Change Distribution Across All Stocks",
        image_path=PLOTS_DIR / 'trump_china_visit_nov_2017_distribution.png',
        image_width=Inches(9)
    )

    # Slide 6: Key Findings
    add_key_findings_slide(prs)

    # Slide 7: Top Gainers
    add_content_slide(
        prs,
        "Top 20 Gaining Stocks During Trump's Visit",
        image_path=PLOTS_DIR / 'trump_china_visit_nov_2017_top_gainers.png',
        image_width=Inches(9)
    )

    # Slide 8: Sector Performance
    add_content_slide(
        prs,
        "Sector Performance Comparison",
        bullet_points=[
            "Semiconductor sector led with +4.8% average gain",
            "Aligned with Qualcomm's $12B deals",
            "",
            "Equipment/Manufacturing: +2.8% average",
            "Benefited from Caterpillar, GE partnerships",
            "",
            "Aviation: +2.4% average",
            "Boeing's $37B aircraft order boost",
            "",
            "Energy: +2.0% average",
            "LNG and shale gas deal impact",
            "",
            "Automotive: +1.3% average",
            "Ford and GM deal benefits"
        ],
        image_path=PLOTS_DIR / 'trump_china_visit_nov_2017_sectors.png',
        image_width=Inches(4.8)
    )

    # Slide 9: Sector Comparison
    add_content_slide(
        prs,
        "Normalized Price Comparison by Sector",
        image_path=PLOTS_DIR / 'sector_comparison_normalized.png',
        image_width=Inches(9)
    )

    # Slide 10-11: Semiconductor Sector
    add_two_image_slide(
        prs,
        "Semiconductor Sector: Qualcomm Deal Beneficiaries",
        PLOTS_DIR / 'stock_300655_semiconductor.png',
        PLOTS_DIR / 'stock_002049_semiconductor.png',
        "300655.SZ 晶瑞电材 (+11.4%)",
        "002049.SZ 紫光国微 (+3.4%)"
    )

    # Slide 12: Energy Sector
    add_two_image_slide(
        prs,
        "Energy Sector: LNG & Shale Gas Deal Impact",
        PLOTS_DIR / 'stock_002221_energy.png',
        PLOTS_DIR / 'stock_601088_energy.png',
        "002221.SZ 东华能源 (+7.4%)",
        "601088.SH 中国神华 (+4.9%)"
    )

    # Slide 13: Aviation Sector
    add_two_image_slide(
        prs,
        "Aviation Sector: Boeing Aircraft Order Impact",
        PLOTS_DIR / 'stock_601021_aviation.png',
        PLOTS_DIR / 'stock_600029_aviation.png',
        "601021.SH 春秋航空 (+6.2%)",
        "600029.SH 南方航空 (+1.2%)"
    )

    # Slide 14: Manufacturing Sector
    add_two_image_slide(
        prs,
        "Manufacturing: Caterpillar & GE Partnership Benefits",
        PLOTS_DIR / 'stock_600031_manufacturing.png',
        PLOTS_DIR / 'stock_000425_manufacturing.png',
        "600031.SH 三一重工 (+7.9%)",
        "000425.SZ 徐工机械 (+3.5%)"
    )

    # Slide 15: Top Gainers Individual
    add_two_image_slide(
        prs,
        "Top Individual Stock Performers",
        PLOTS_DIR / 'stock_002409_top gainer.png',
        PLOTS_DIR / 'stock_603501_top gainer.png',
        "002409.SZ 雅克科技 (+33.1%)",
        "603501.SH 豪威集团 (+28.2%)"
    )

    # Slide 16: Price vs Volume
    add_content_slide(
        prs,
        "Price Change vs Volume Change Analysis",
        bullet_points=[
            "High volume + price increase indicates strong buying interest",
            "Many gainers showed increased trading volume",
            "Top performers attracted significant investor attention",
            "Volume surge validates market confidence in trade deals"
        ],
        image_path=PLOTS_DIR / 'trump_china_visit_nov_2017_price_volume.png',
        image_width=Inches(5)
    )

    # Slide 17: Conclusions
    add_content_slide(
        prs,
        "Conclusions",
        bullet_points=[
            "Trump's China visit had a measurable positive impact on Chinese stocks",
            "",
            "Sector-specific gains aligned with announced trade deals:",
            "   - Semiconductor stocks rose on Qualcomm smartphone deals",
            "   - Energy stocks benefited from LNG/shale gas investments",
            "   - Aviation stocks gained from Boeing aircraft orders",
            "   - Manufacturing stocks rose on equipment partnerships",
            "",
            "Market showed broad-based optimism with 61.7% of stocks gaining",
            "",
            "The visit demonstrated how diplomatic events can create",
            "significant short-term market opportunities",
            "",
            "Investors who anticipated deal announcements could have",
            "positioned for sector-specific gains"
        ]
    )

    # Slide 18: Methodology
    add_content_slide(
        prs,
        "Methodology & Data Sources",
        bullet_points=[
            "Data Source: Tushare Pro API (daily OHLCV data)",
            "",
            "Analysis Period:",
            "   - Pre-event: October 30 - November 7, 2017",
            "   - Event: November 8-10, 2017 (Trump's visit)",
            "   - Post-event: November 13-30, 2017",
            "",
            "Stocks Analyzed: 2,261 A-share stocks (SH + SZ)",
            "",
            "Metrics:",
            "   - Price change = (Event close - Pre-event close) / Pre-event close",
            "   - Volume change = (Event avg volume - Pre-event avg volume) / Pre-event avg",
            "",
            "Sector Classification: Based on company names and deal relevance",
            "",
            "Analysis Tool: Python with pandas, matplotlib, Tushare"
        ]
    )

    # Save presentation
    output_path = OUTPUT_DIR / 'Trump_China_Visit_Stock_Analysis.pptx'
    prs.save(str(output_path))
    print(f"\nPresentation saved to: {output_path}")

    return output_path


if __name__ == '__main__':
    create_presentation()
